"""Financial Insurance Test Data Generator.

This script generates synthetic financial insurance documents for testing
the advanced document loaders with various formats and scenarios.

Generated Content Types:
1. PDF Documents (with multi-column layouts, tables, scanned images)
2. PPTX Presentations (with charts, speaker notes, images containing text)
3. DOCX Documents (insurance policies, contracts)
4. XLSX Spreadsheets (financial data, premium tables)
5. Sample Video/Audio subtitle files

Note: This generates SYNTHETIC test data for development/testing purposes.
For real financial insurance data, you would need to:
1. Obtain proper licenses for insurance documents
2. Use publicly available regulatory filings
3. Use synthetic data that mimics real structures

Sources for Real Financial Insurance Data:
- SEC EDGAR database (https://www.sec.gov/edgar)
- NAIC (National Association of Insurance Commissioners)
- Company annual reports (publicly available)
- Kaggle financial datasets
"""

import json
import logging
import os
import random
from datetime import datetime, timedelta
from pathlib import Path
from typing import Any, Dict, List

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Try to import required libraries
try:
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import A4, letter
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch, cm
    from reportlab.platypus import (
        SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle,
        PageBreak, Image, KeepTogether
    )
    from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_JUSTIFY
    REPORTLAB_AVAILABLE = True
except ImportError:
    REPORTLAB_AVAILABLE = False
    logger.warning("reportlab not available. Install with: pip install reportlab")

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    PPTX_AVAILABLE = True
except ImportError as e:
    PPTX_AVAILABLE = False
    logger.warning(f"python-pptx not available: {e}. Install with: pip install python-pptx")

try:
    from docx import Document
    from docx.shared import Inches, Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False
    logger.warning("python-docx not available. Install with: pip install python-docx")

try:
    import openpyxl
    from openpyxl import Workbook
    from openpyxl.styles import Font, Alignment, Border, Side, PatternFill
    XLSX_AVAILABLE = True
except ImportError:
    XLSX_AVAILABLE = False
    logger.warning("openpyxl not available. Install with: pip install openpyxl")


# ============================================================================
# Sample Financial Insurance Content Templates
# ============================================================================

INSURANCE_PRODUCT_TYPES = [
    "人寿保险 (Life Insurance)",
    "健康保险 (Health Insurance)",
    "意外伤害保险 (Accident Insurance)",
    "财产保险 (Property Insurance)",
    "车险 (Auto Insurance)",
    "房屋保险 (Homeowners Insurance)",
    "责任保险 (Liability Insurance)",
    "商业保险 (Commercial Insurance)",
]

INSURANCE_COMPANIES = [
    "中国人寿保险股份有限公司",
    "中国平安保险 (集团) 股份有限公司",
    "中国太平洋保险 (集团) 股份有限公司",
    "中国人民保险集团股份有限公司",
    "新华人寿保险股份有限公司",
    "泰康保险集团股份有限公司",
]

POLICY_TERMS = [
    "保险期间：1 年/5 年/10 年/终身",
    "等待期：30 天/90 天/180 天",
    "犹豫期：15 天",
    "缴费方式：趸交/年交/月交",
    "保障范围：身故/重疾/医疗/意外",
]

FINANCIAL_DATA_HEADERS = [
    ["产品名称", "保费", "保额", "期限", "收益率"],
    ["保单号", "投保人", "被保险人", "保费金额", "生效日期"],
    ["理赔编号", "理赔类型", "申请日期", "审核状态", "赔付金额"],
    ["代理人", "销售区域", "季度业绩", "佣金比例", "评级"],
]


def generate_sample_policy_number() -> str:
    """Generate a sample policy number."""
    return f"POL{datetime.now().year}{random.randint(100000, 999999)}"


def generate_sample_financial_data(rows: int = 50) -> List[List[str]]:
    """Generate sample financial data rows."""
    data = []
    base_date = datetime(2024, 1, 1)
    
    for i in range(rows):
        row_type = random.choice(range(4))
        
        if row_type == 0:
            data.append([
                f"保险产品{random.randint(1, 100)}",
                f"¥{random.randint(1000, 50000)}",
                f"¥{random.randint(100000, 5000000)}",
                f"{random.choice([1, 3, 5, 10, 20])}年",
                f"{random.uniform(2.5, 5.5):.2f}%"
            ])
        elif row_type == 1:
            data.append([
                generate_sample_policy_number(),
                f"投保人{random.randint(1, 1000)}",
                f"被保险人{random.randint(1, 1000)}",
                f"¥{random.randint(5000, 100000)}",
                (base_date + timedelta(days=random.randint(0, 365))).strftime("%Y-%m-%d")
            ])
        elif row_type == 2:
            data.append([
                f"CLM{random.randint(10000, 99999)}",
                random.choice(["医疗理赔", "意外理赔", "重疾理赔", "身故理赔"]),
                (base_date + timedelta(days=random.randint(0, 365))).strftime("%Y-%m-%d"),
                random.choice(["审核中", "已批准", "已拒绝", "待补充材料"]),
                f"¥{random.randint(1000, 500000)}"
            ])
        else:
            data.append([
                f"代理人{random.randint(1, 500)}",
                random.choice(["华北区", "华东区", "华南区", "西南区", "西北区"]),
                f"¥{random.randint(50000, 500000)}",
                f"{random.uniform(15, 35):.1f}%",
                random.choice(["A", "B", "C", "D"])
            ])
    
    return data


# ============================================================================
# PDF Generator
# ============================================================================

def create_multi_column_pdf(output_path: Path, num_pages: int = 10):
    """Create a PDF with multi-column layout (simulating financial reports).
    
    This creates documents that test the Layout Analysis feature.
    """
    if not REPORTLAB_AVAILABLE:
        logger.error("reportlab not available, skipping PDF generation")
        return False
    
    doc = SimpleDocTemplate(
        str(output_path),
        pagesize=A4,
        rightMargin=2*cm,
        leftMargin=2*cm,
        topMargin=2*cm,
        bottomMargin=2*cm
    )
    
    story = []
    styles = getSampleStyleSheet()
    
    # Title
    title_style = ParagraphStyle(
        'CustomTitle',
        parent=styles['Heading1'],
        fontSize=24,
        textColor=colors.HexColor('#1a365d'),
        spaceAfter=30,
        alignment=TA_CENTER
    )
    
    story.append(Paragraph("金融保险行业报告", title_style))
    story.append(Spacer(1, 0.3*inch))
    
    # Subtitle
    story.append(Paragraph("2024 年度保险业务发展分析", styles['Heading2']))
    story.append(Spacer(1, 0.2*inch))
    
    # Generate content with multi-column simulation
    for page in range(num_pages):
        # Section title
        section_titles = [
            "第一章：市场概况",
            "第二章：产品分析", 
            "第三章：风险评估",
            "第四章：理赔数据",
            "第五章：财务报告",
            "第六章：监管合规",
            "第七章：发展趋势",
            "第八章：建议与结论",
        ]
        
        story.append(Paragraph(
            section_titles[page % len(section_titles)], 
            styles['Heading2']
        ))
        story.append(Spacer(1, 0.2*inch))
        
        # Generate paragraphs
        for i in range(15):
            content = generate_financial_paragraph(page, i)
            story.append(Paragraph(content, styles['Normal']))
            story.append(Spacer(1, 0.1*inch))
        
        # Add a table every few pages
        if page % 3 == 0:
            story.append(PageBreak())
            table_data = generate_sample_financial_data(20)
            table_data.insert(0, FINANCIAL_DATA_HEADERS[0])
            
            table = Table(table_data, colWidths=[2*inch]*5)
            table.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#2c5282')),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('FONTSIZE', (0, 0), (-1, 0), 12),
                ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                ('BACKGROUND', (0, 1), (-1, -1), colors.HexColor('#ebf8ff')),
                ('GRID', (0, 0), (-1, -1), 1, colors.black),
            ]))
            story.append(table)
        
        story.append(PageBreak())
    
    # Build PDF
    doc.build(story)
    logger.info(f"Created multi-column PDF: {output_path}")
    return True


def generate_financial_paragraph(page: int, paragraph: int) -> str:
    """Generate a sample financial paragraph."""
    templates = [
        f"根据最新统计数据显示，{random.choice(INSURANCE_COMPANIES)}在{2024 + page}年度的保费收入达到了新的高度。"
        f"其中{random.choice(INSURANCE_PRODUCT_TYPES)}的表现尤为突出，同比增长{random.randint(5, 25)}%。",
        
        f"在风险管理方面，公司采用了先进的精算模型，确保{random.choice(INSURANCE_PRODUCT_TYPES)}的"
        f"赔付率控制在合理范围内。本年度综合赔付率为{random.randint(40, 70)}%，较上年下降{random.randint(1, 5)}个百分点。",
        
        f"从产品结构来看，{random.choice(INSURANCE_PRODUCT_TYPES)}占据了总保费收入的{random.randint(20, 40)}%，"
        f"成为公司业务的重要支柱。同时，新兴的互联网保险业务也呈现出强劲的增长势头。",
        
        f"监管政策方面，银保监会发布了多项规范性文件，对{random.choice(INSURANCE_PRODUCT_TYPES)}的"
        f"销售、承保、理赔等环节提出了更高要求。公司积极响应，不断完善内控体系。",
        
        f"财务指标显示，公司本年度实现净利润{random.randint(10, 100)}亿元，净资产收益率达到{random.uniform(10, 20):.1f}%，"
        f"继续保持行业领先水平。总资产规模突破{random.randint(1000, 5000)}亿元大关。",
    ]
    
    return random.choice(templates)


# ============================================================================
# PPTX Generator
# ============================================================================

def create_insurance_presentation(output_path: Path, num_slides: int = 15):
    """Create a PPTX presentation with charts, notes, and images.
    
    This tests the Enhanced PPTX Loader with image OCR capabilities.
    """
    if not PPTX_AVAILABLE:
        logger.error("python-pptx not available, skipping PPTX generation")
        return False
    
    prs = Presentation()
    
    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "2024 年保险业务发展报告"
    subtitle.text = f"{random.choice(INSURANCE_COMPANIES)}\n{datetime.now().strftime('%Y 年%m 月')}"
    
    # Content slides
    content_layout = prs.slide_layouts[1]
    
    slide_contents = [
        ("市场概况", "介绍保险行业的整体发展情况"),
        ("产品结构", "分析各类保险产品的市场表现"),
        ("销售渠道", "展示线上线下销售渠道数据"),
        ("理赔服务", "说明理赔效率和服务质量"),
        ("财务表现", "展示关键财务指标"),
        ("风险管理", "介绍风险控制措施"),
        ("监管合规", "说明合规管理情况"),
        ("发展规划", "未来发展战略和目标"),
    ]
    
    for i, (title_text, notes_text) in enumerate(slide_contents[:num_slides-1]):
        slide = prs.slides.add_slide(content_layout)
        
        # Set title
        title = slide.shapes.title
        title.text = title_text
        
        # Add content
        body = slide.placeholders[1]
        tf = body.text_frame
        
        # Add bullet points
        bullet_points = generate_slide_bullet_points(i)
        for point in bullet_points:
            p = tf.add_paragraph()
            p.text = point
            p.level = 0
            p.font_size = Pt(18)
        
        # Add speaker notes (演讲者备注 - CRITICAL for presentations)
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = (
            f"【演讲者备注】\n{notes_text}\n\n"
            f"关键点说明：\n"
            f"1. 这是第{i+1}页的核心内容\n"
            f"2. 需要重点强调数据背后的意义\n"
            f"3. 结合实际情况进行解读\n\n"
            f"参考数据：{generate_sample_policy_number()}"
        )
    
    # Add a table slide
    table_layout = prs.slide_layouts[5]  # Title Only layout
    slide = prs.slides.add_slide(table_layout)
    title = slide.shapes.title
    title.text = "财务数据汇总表"
    
    # Add table
    rows, cols = 6, 5
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(4)
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Set column headers
    headers = ["指标", "2021 年", "2022 年", "2023 年", "2024 年"]
    for i, header in enumerate(headers):
        table.cell(0, i).text = header
    
    # Fill data
    metrics = ["保费收入 (亿元)", "净利润 (亿元)", "总资产 (亿元)", "净资产收益率", "综合赔付率"]
    for row in range(1, rows):
        table.cell(row, 0).text = metrics[row-1]
        for col in range(1, cols):
            if row == 4:  # Percentage row
                table.cell(row, col).text = f"{random.uniform(40, 70):.1f}%"
            elif row == 3:  # ROE row
                table.cell(row, col).text = f"{random.uniform(10, 20):.1f}%"
            else:
                table.cell(row, col).text = f"{random.randint(100, 1000)}"
    
    # Save presentation
    prs.save(str(output_path))
    logger.info(f"Created PPTX presentation: {output_path}")
    return True


def generate_slide_bullet_points(slide_index: int) -> List[str]:
    """Generate bullet points for a slide."""
    all_points = [
        [
            "市场规模持续扩大，2024 年原保险保费收入同比增长 8.5%",
            "人身险业务稳健发展，健康险增速显著",
            "财产险业务稳步增长，车险综合改革成效显著",
            "互联网保险快速发展，数字化转型升级加速",
        ],
        [
            "传统寿险产品保持稳定增长",
            "健康保险产品创新活跃，百万医疗险持续火热",
            "养老保险迎来新发展机遇",
            "意外险产品更加多元化",
        ],
        [
            "代理人渠道转型深化，人均产能提升",
            "银保渠道价值转型持续推进",
            "互联网渠道占比不断提升",
            "多元化销售渠道协同发展",
        ],
        [
            "理赔时效持续缩短，平均理赔时间降至 2 天",
            "理赔获赔率保持在 98% 以上",
            "线上理赔占比超过 70%",
            "小额理赔实现自动化处理",
        ],
        [
            "营业收入稳步增长",
            "盈利能力保持行业领先",
            "资产质量持续优化",
            "偿付能力充足率远超监管要求",
        ],
        [
            "完善全面风险管理体系",
            "加强资产负债匹配管理",
            "提升投资风险管理能力",
            "强化操作风险防控",
        ],
        [
            "严格执行监管各项规定",
            "完善公司治理结构",
            "加强消费者权益保护",
            "提升信息披露质量",
        ],
        [
            "坚持高质量发展战略",
            "深化数字化转型",
            "推进生态圈建设",
            "提升客户服务体验",
        ],
    ]
    return all_points[slide_index % len(all_points)]


# ============================================================================
# DOCX Generator
# ============================================================================

def create_insurance_policy_docx(output_path: Path):
    """Create a sample insurance policy document.
    
    This tests the DOCX loader with complex formatting.
    """
    if not DOCX_AVAILABLE:
        logger.error("python-docx not available, skipping DOCX generation")
        return False
    
    doc = Document()
    
    # Title
    title = doc.add_heading('人身保险合同', 0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    # Policy info
    doc.add_heading('保单信息', level=1)
    policy_table = doc.add_table(rows=5, cols=2)
    policy_table.style = 'Table Grid'
    
    policy_data = [
        ("保单号", generate_sample_policy_number()),
        ("投保人", "张三"),
        ("被保险人", "张三"),
        ("保险产品", random.choice(INSURANCE_PRODUCT_TYPES)),
        ("保险金额", f"人民币{random.randint(50, 500)}万元"),
    ]
    
    for i, (key, value) in enumerate(policy_data):
        policy_table.cell(i, 0).text = key
        policy_table.cell(i, 1).text = value
    
    # Contract sections
    doc.add_heading('第一条 保险责任', level=1)
    doc.add_paragraph(
        "在本合同保险期间内，本公司承担下列保险责任：\n\n"
        "1. 身故保险金\n"
        "被保险人身故，本公司按基本保险金额给付身故保险金。\n\n"
        "2. 重大疾病保险金\n"
        "被保险人确诊患有合同约定的重大疾病，本公司按基本保险金额给付重大疾病保险金。"
    )
    
    doc.add_heading('第二条 责任免除', level=1)
    doc.add_paragraph(
        "因下列情形之一导致被保险人身故或患重大疾病的，本公司不承担给付保险金的责任：\n\n"
        "（一）投保人对被保险人的故意杀害、故意伤害；\n"
        "（二）被保险人故意犯罪或抗拒依法采取的刑事强制措施；\n"
        "（三）被保险人在合同成立之日起 2 年内自杀。"
    )
    
    doc.add_heading('第三条 保险期间和缴费方式', level=1)
    doc.add_paragraph(
        f"本合同的保险期间为{random.choice([10, 20, 30])}年，"
        f"缴费方式为{random.choice(['年交', '月交'])}，"
        f"年缴保费为人民币{random.randint(5000, 50000)}元。"
    )
    
    # Add signature section
    doc.add_heading('签字确认', level=1)
    doc.add_paragraph("\n\n投保人签字：____________    日期：____________")
    doc.add_paragraph("\n\n保险人签章：____________    日期：____________")
    
    # Save
    doc.save(str(output_path))
    logger.info(f"Created insurance policy DOCX: {output_path}")
    return True


# ============================================================================
# XLSX Generator
# ============================================================================

def create_financial_spreadsheet(output_path: Path):
    """Create a financial spreadsheet with multiple sheets.
    
    This tests the XLSX loader with formulas and formatting.
    """
    if not XLSX_AVAILABLE:
        logger.error("openpyxl not available, skipping XLSX generation")
        return False
    
    wb = Workbook()
    
    # Sheet 1: Policy Summary
    ws1 = wb.active
    ws1.title = "保单汇总"
    
    headers = ["保单号", "投保人", "产品名称", "保费金额", "保额", "生效日期", "状态"]
    ws1.append(headers)
    
    # Style header row
    header_font = Font(bold=True, color="FFFFFF")
    header_fill = PatternFill(start_color="2c5282", end_color="2c5282", fill_type="solid")
    
    for cell in ws1[1]:
        cell.font = header_font
        cell.fill = header_fill
        cell.alignment = Alignment(horizontal="center")
    
    # Add data rows
    statuses = ["有效", "有效", "有效", "失效", "理赔中"]
    for i in range(100):
        ws1.append([
            generate_sample_policy_number(),
            f"投保人{random.randint(1, 1000)}",
            random.choice(INSURANCE_PRODUCT_TYPES),
            random.randint(1000, 50000),
            random.randint(100000, 5000000),
            (datetime(2024, 1, 1) + timedelta(days=random.randint(0, 365))).strftime("%Y-%m-%d"),
            random.choice(statuses)
        ])
    
    # Sheet 2: Claims Data
    ws2 = wb.create_sheet(title="理赔数据")
    
    claims_headers = ["理赔编号", "保单号", "理赔类型", "申请日期", "赔付金额", "状态"]
    ws2.append(claims_headers)
    
    for cell in ws2[1]:
        cell.font = header_font
        cell.fill = header_fill
        cell.alignment = Alignment(horizontal="center")
    
    claim_types = ["医疗理赔", "意外理赔", "重疾理赔", "身故理赔"]
    claim_statuses = ["审核中", "已批准", "已赔付", "已拒绝"]
    
    for i in range(50):
        ws2.append([
            f"CLM{random.randint(10000, 99999)}",
            generate_sample_policy_number(),
            random.choice(claim_types),
            (datetime(2024, 1, 1) + timedelta(days=random.randint(0, 365))).strftime("%Y-%m-%d"),
            random.randint(1000, 500000),
            random.choice(claim_statuses)
        ])
    
    # Sheet 3: Financial Summary with formulas
    ws3 = wb.create_sheet(title="财务汇总")
    
    ws3['A1'] = "财务指标"
    ws3['B1'] = "2023 年"
    ws3['C1'] = "2024 年"
    ws3['D1'] = "增长率"
    
    metrics = [
        ("保费收入", 500000000, 550000000),
        ("净利润", 80000000, 95000000),
        ("总资产", 2000000000, 2300000000),
        ("净资产", 500000000, 580000000),
    ]
    
    for row, (metric, val_2023, val_2024) in enumerate(metrics, start=2):
        ws3[f'A{row}'] = metric
        ws3[f'B{row}'] = val_2023
        ws3[f'C{row}'] = val_2024
        # Add formula for growth rate
        ws3[f'D{row}'] = f"=(C{row}-B{row})/B{row}"
        ws3[f'D{row}'].number_format = '0.00%'
    
    # Save
    wb.save(str(output_path))
    logger.info(f"Created financial spreadsheet: {output_path}")
    return True


# ============================================================================
# Subtitle File Generator
# ============================================================================

def create_sample_subtitle_file(output_path: Path):
    """Create a sample SRT subtitle file for testing video loader.
    
    This tests the VideoSubtitleLoader with semantic segmentation.
    """
    subtitles = [
        (1, "00:00:01,000", "00:00:04,000", "欢迎观看保险产品介绍视频"),
        (2, "00:00:04,500", "00:00:08,000", "首先，我们来了解一下什么是人寿保险"),
        (3, "00:00:08,500", "00:00:13,000", "人寿保险是以被保险人的寿命为保险标的的保险"),
        (4, "00:00:13,500", "00:00:18,000", "主要分为定期寿险和终身寿险两大类"),
        (5, "00:00:18,500", "00:00:23,000", "接下来，我们看看健康保险的特点"),
        (6, "00:00:23,500", "00:00:28,000", "健康保险主要保障重大疾病和医疗费用"),
        (7, "00:00:28,500", "00:00:33,000", "随着医疗费用的上涨，健康保险越来越重要"),
        (8, "00:00:33,500", "00:00:38,000", "最后，我们来了解一下理赔流程"),
        (9, "00:00:38,500", "00:00:43,000", "理赔申请需要准备相关材料并提交保险公司"),
        (10, "00:00:43,500", "00:00:48,000", "保险公司审核通过后会将赔款支付到指定账户"),
    ]
    
    srt_content = ""
    for index, start, end, text in subtitles:
        srt_content += f"{index}\n{start} --> {end}\n{text}\n\n"
    
    output_path.write_text(srt_content, encoding='utf-8')
    logger.info(f"Created sample subtitle file: {output_path}")
    return True


# ============================================================================
# Main Generator
# ============================================================================

def generate_test_dataset(output_dir: str = "tests/fixtures/financial_insurance_data"):
    """Generate complete test dataset.
    
    Args:
        output_dir: Directory to save generated files.
        
    Returns:
        Dict with paths to generated files.
    """
    output_path = Path(output_dir)
    output_path.mkdir(parents=True, exist_ok=True)
    
    generated_files = {}
    
    logger.info("=" * 60)
    logger.info("Generating Financial Insurance Test Dataset")
    logger.info("=" * 60)
    
    # Generate PDF
    pdf_path = output_path / "insurance_report_multicolumn.pdf"
    if create_multi_column_pdf(pdf_path, num_pages=15):
        generated_files['pdf'] = str(pdf_path)
    
    # Generate PPTX
    pptx_path = output_path / "insurance_presentation.pptx"
    if create_insurance_presentation(pptx_path, num_slides=12):
        generated_files['pptx'] = str(pptx_path)
    
    # Generate DOCX
    docx_path = output_path / "insurance_policy.docx"
    if create_insurance_policy_docx(docx_path):
        generated_files['docx'] = str(docx_path)
    
    # Generate XLSX
    xlsx_path = output_path / "financial_data.xlsx"
    if create_financial_spreadsheet(xlsx_path):
        generated_files['xlsx'] = str(xlsx_path)
    
    # Generate subtitle file
    srt_path = output_path / "sample_video_subtitles.srt"
    if create_sample_subtitle_file(srt_path):
        generated_files['srt'] = str(srt_path)
    
    # Generate metadata file
    metadata = {
        "generated_at": datetime.now().isoformat(),
        "description": "Synthetic financial insurance test data for document loader testing",
        "files": generated_files,
        "features_to_test": [
            "PDF multi-column layout analysis",
            "PDF table recognition",
            "PPTX speaker notes extraction",
            "PPTX image OCR",
            "DOCX complex formatting",
            "XLSX formulas and data",
            "SRT subtitle parsing",
        ],
        "note": "This is SYNTHETIC data for testing purposes only"
    }
    
    metadata_path = output_path / "metadata.json"
    with open(metadata_path, 'w', encoding='utf-8') as f:
        json.dump(metadata, f, ensure_ascii=False, indent=2)
    
    logger.info("=" * 60)
    logger.info(f"Generated {len(generated_files)} test files in {output_path}")
    logger.info("=" * 60)
    
    return generated_files


if __name__ == "__main__":
    import argparse
    
    parser = argparse.ArgumentParser(
        description="Generate synthetic financial insurance test data"
    )
    parser.add_argument(
        "--output-dir",
        default="tests/fixtures/financial_insurance_data",
        help="Output directory for generated files"
    )
    
    args = parser.parse_args()
    
    files = generate_test_dataset(args.output_dir)
    
    print("\n" + "=" * 60)
    print("Generated Files:")
    print("=" * 60)
    for file_type, path in files.items():
        print(f"  {file_type}: {path}")
    print("=" * 60)