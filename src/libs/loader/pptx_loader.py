"""PPTX Loader with semantic-aware parsing.

This module implements PPTX parsing with preservation of:
- Slide structure and hierarchy
- Speaker notes (演讲者备注) - often more critical than content
- Master slides and layouts
- Embedded images and media
- Shapes and text boxes
- Transitions and animations metadata
- Hidden slides

Key Features:
- Extracts speaker notes as priority content
- Preserves slide hierarchy with metadata
- Converts shapes and diagrams to structured descriptions
- Tracks presentation flow and narrative
"""

from __future__ import annotations

import hashlib
import logging
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

try:
    from pptx import Presentation
    from pptx.slide import Slide
    from pptx.shapes.base import BaseShape
    from pptx.shapes.picture import Picture
    from pptx.shapes.autoshape import Shape as AutoShape
    from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
    # Note: Placeholder class import path changed in python-pptx 1.0+
    # We use placeholder_format.type attribute instead of importing Placeholder
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

from src.core.types import Document
from src.libs.loader.base_loader import BaseLoader

logger = logging.getLogger(__name__)


class PptxLoader(BaseLoader):
    """PPTX Loader with semantic-aware parsing.
    
    This loader preserves:
    1. Slide structure: titles, content, hierarchy
    2. Speaker Notes (演讲者备注): Extracted as priority content
    3. Master slides: Layout and theme information
    4. Images: Extracted with slide context
    5. Shapes: Text boxes, diagrams converted to markdown
    6. Hidden slides: Marked separately
    7. Transitions/Animations: Metadata tracking
    
    Configuration:
        extract_notes: Extract speaker notes (default: True, HIGH PRIORITY)
        extract_masters: Extract master slide info (default: True)
        extract_images: Extract embedded images (default: True)
        extract_shapes: Extract shape text content (default: True)
        include_hidden: Include hidden slides (default: True)
        image_storage_dir: Base directory for image storage
    """
    
    def __init__(
        self,
        extract_notes: bool = True,
        extract_masters: bool = True,
        extract_images: bool = True,
        extract_shapes: bool = True,
        include_hidden: bool = True,
        image_storage_dir: str | Path = "data/images",
    ):
        """Initialize PPTX Loader.
        
        Args:
            extract_notes: Whether to extract speaker notes (CRITICAL).
            extract_masters: Whether to extract master slide info.
            extract_images: Whether to extract embedded images.
            extract_shapes: Whether to extract shape text content.
            include_hidden: Whether to include hidden slides.
            image_storage_dir: Directory for storing extracted images.
        """
        if not PPTX_AVAILABLE:
            raise ImportError(
                "python-pptx is required for PptxLoader. "
                "Install with: pip install python-pptx"
            )
        
        self.extract_notes = extract_notes
        self.extract_masters = extract_masters
        self.extract_images = extract_images
        self.extract_shapes = extract_shapes
        self.include_hidden = include_hidden
        self.image_storage_dir = Path(image_storage_dir)
    
    def load(self, file_path: str | Path) -> Document:
        """Load and parse a PPTX file.
        
        Args:
            file_path: Path to the PPTX file.
            
        Returns:
            Document with Markdown text and metadata.
            
        Raises:
            FileNotFoundError: If the PPTX file doesn't exist.
            ValueError: If the file is not a valid PPTX.
            RuntimeError: If parsing fails.
        """
        path = self._validate_file(file_path)
        if path.suffix.lower() != '.pptx':
            raise ValueError(f"File is not a PPTX: {path}")
        
        # Compute document hash
        doc_hash = self._compute_file_hash(path)
        doc_id = f"doc_{doc_hash[:16]}"
        
        # Parse PPTX
        prs = Presentation(path)
        
        # Extract metadata
        metadata = self._extract_metadata(prs, path, doc_hash)
        
        # Extract master slide info
        master_info = self._extract_master_slides(prs) if self.extract_masters else ""
        
        # Extract slides content
        slides_content, slide_metadata = self._extract_slides(prs, doc_hash)
        
        # Build final markdown content
        sections = []
        
        # Title slide / presentation info
        if metadata.get('title'):
            sections.append(f"# {metadata['title']}\n")
        
        # Master slide info (optional, for context)
        if master_info.strip():
            sections.append(master_info)
        
        # Main slides content
        sections.append(slides_content)
        
        # Combine sections
        text_content = "\n\n".join(sections)
        
        # Merge slide-level metadata
        metadata['slides'] = slide_metadata
        metadata['total_slides'] = len(prs.slides)
        metadata['slides_with_notes'] = sum(1 for s in slide_metadata if s.get('has_notes'))
        metadata['hidden_slides'] = sum(1 for s in slide_metadata if s.get('is_hidden'))
        
        return Document(
            id=doc_id,
            text=text_content,
            metadata=metadata
        )
    
    def _extract_metadata(
        self, 
        prs: Presentation, 
        path: Path, 
        doc_hash: str
    ) -> Dict[str, Any]:
        """Extract presentation metadata.
        
        Args:
            prs: Parsed presentation object.
            path: File path.
            doc_hash: Document hash.
            
        Returns:
            Metadata dictionary.
        """
        metadata = {
            "source_path": str(path),
            "doc_type": "pptx",
            "doc_hash": doc_hash,
        }
        
        # Core properties
        core_props = prs.core_properties
        if core_props:
            if core_props.title:
                metadata["title"] = core_props.title
            if core_props.author:
                metadata["author"] = core_props.author
            if core_props.created:
                metadata["created_date"] = core_props.created.isoformat()
            if core_props.modified:
                metadata["modified_date"] = core_props.modified.isoformat()
            if core_props.subject:
                metadata["subject"] = core_props.subject
            if core_props.keywords:
                metadata["keywords"] = core_props.keywords
        
        # Presentation info
        metadata["slide_count"] = len(prs.slides)
        
        return metadata
    
    def _extract_master_slides(self, prs: Presentation) -> str:
        """Extract master slide information.
        
        Args:
            prs: Parsed presentation object.
            
        Returns:
            Markdown-formatted master slide info.
        """
        lines = ["\n## 幻灯片母版信息 (Master Slides)\n"]
        
        try:
            for idx, master in enumerate(prs.slide_masters):
                lines.append(f"\n### 母版 {idx + 1}: {master.name}")
                
                # Get layout names
                layouts = master.slide_layouts
                if layouts:
                    lines.append(f"\n**布局数量:** {len(layouts)}")
                    lines.append("\n**可用布局:**")
                    for layout in layouts[:10]:  # Limit to first 10
                        lines.append(f"- {layout.name}")
                    if len(layouts) > 10:
                        lines.append(f"- ... 共 {len(layouts)} 个布局")
                        
        except Exception as e:
            logger.debug(f"Could not extract master slides: {e}")
        
        return "\n".join(lines)
    
    def _extract_slides(
        self, 
        prs: Presentation, 
        doc_hash: str
    ) -> Tuple[str, List[Dict[str, Any]]]:
        """Extract all slides content.
        
        Args:
            prs: Parsed presentation object.
            doc_hash: Document hash for image paths.
            
        Returns:
            Tuple of (markdown_content, slide_metadata_list)
        """
        markdown_parts = []
        slide_metadata_list = []
        image_index = 0
        
        for idx, slide in enumerate(prs.slides):
            # Check if hidden
            is_hidden = getattr(slide, '_element', None) is not None and \
                       slide._element.get('{http://schemas.openxmlformats.org/presentationml/2006/main}hide') == '1'
            
            if is_hidden and not self.include_hidden:
                continue
            
            # Slide header
            slide_header = f"\n\n## 幻灯片 {idx + 1}\n"
            markdown_parts.append(slide_header)
            
            slide_meta = {
                "slide_number": idx + 1,
                "is_hidden": is_hidden,
                "has_notes": False,
                "has_images": False,
                "shape_count": len(slide.shapes),
            }
            
            # Extract title
            title = self._extract_slide_title(slide)
            if title:
                markdown_parts.append(f"### {title}\n")
                slide_meta["title"] = title
            
            # Extract shapes content
            shapes_content, shapes_meta = self._extract_shapes(slide, doc_hash, idx, image_index)
            if shapes_content.strip():
                markdown_parts.append(shapes_content)
                image_index = shapes_meta.get('image_index', image_index)
                if shapes_meta.get('has_images'):
                    slide_meta["has_images"] = True
                    slide_meta["image_count"] = shapes_meta.get('image_count', 0)
            
            # Extract speaker notes (CRITICAL - often more important than slide content)
            if self.extract_notes:
                notes = self._extract_notes(slide)
                if notes:
                    markdown_parts.append(f"\n\n**【演讲者备注】**:\n> {notes}\n")
                    slide_meta["has_notes"] = True
                    slide_meta["notes_preview"] = notes[:200]  # Preview for search
            
            slide_metadata_list.append(slide_meta)
        
        return "\n".join(markdown_parts), slide_metadata_list
    
    def _extract_slide_title(self, slide: Slide) -> Optional[str]:
        """Extract title from slide.
        
        Args:
            slide: Slide object.
            
        Returns:
            Title string or None.
        """
        # Try to get title from title placeholder
        for shape in slide.shapes:
            if shape.is_placeholder:
                ph_type = shape.placeholder_format.type
                if ph_type == PP_PLACEHOLDER.TITLE or ph_type == PP_PLACEHOLDER.CENTER_TITLE:
                    if shape.has_text_frame and shape.text_frame.text.strip():
                        return shape.text_frame.text.strip()
        
        # Fallback: first shape with text
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                text = shape.text_frame.text.strip().split('\n')[0]
                if len(text) < 100:  # Reasonable title length
                    return text
        
        return None
    
    def _extract_shapes(
        self, 
        slide: Slide, 
        doc_hash: str, 
        slide_idx: int,
        start_image_idx: int
    ) -> Tuple[str, Dict[str, Any]]:
        """Extract content from shapes on a slide.
        
        Args:
            slide: Slide object.
            doc_hash: Document hash.
            slide_idx: Slide index.
            start_image_idx: Starting image index.
            
        Returns:
            Tuple of (markdown_content, metadata)
        """
        markdown_parts = []
        image_idx = start_image_idx
        has_images = False
        
        for shape_idx, shape in enumerate(slide.shapes):
            # Skip title placeholder (already extracted)
            if shape.is_placeholder and shape.placeholder_format.type == PP_PLACEHOLDER.TITLE:
                continue
            
            # Handle text boxes / auto shapes
            if self.extract_shapes and shape.has_text_frame:
                text_content = self._shape_to_markdown(shape)
                if text_content.strip():
                    markdown_parts.append(text_content)
            
            # Handle images
            if self.extract_images and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_meta = self._extract_image(shape, doc_hash, slide_idx, image_idx)
                if image_meta:
                    has_images = True
                    placeholder = f"\n[IMAGE: {image_meta['id']}]\n"
                    markdown_parts.append(placeholder)
                    image_idx += 1
            
            # Handle tables
            if shape.has_table:
                table_md = self._table_to_markdown(shape.table)
                markdown_parts.append(table_md)
        
        metadata = {
            'image_index': image_idx,
            'has_images': has_images,
            'image_count': image_idx - start_image_idx
        }
        
        return "\n\n".join(markdown_parts), metadata
    
    def _shape_to_markdown(self, shape: BaseShape) -> str:
        """Convert a shape to Markdown.
        
        Args:
            shape: Shape object.
            
        Returns:
            Markdown-formatted text.
        """
        if not shape.has_text_frame:
            return ""
        
        text_frame = shape.text_frame
        parts = []
        
        for para_idx, para in enumerate(text_frame.paragraphs):
            para_text = ""
            
            for run in para.runs:
                run_text = run.text
                if not run_text:
                    continue
                
                # Apply formatting
                if run.font.bold and run.font.italic:
                    run_text = f"***{run_text}***"
                elif run.font.bold:
                    run_text = f"**{run_text}**"
                elif run.font.italic:
                    run_text = f"*{run_text}*"
                
                para_text += run_text
            
            if para_text.strip():
                # Determine indentation level
                indent = "  " * (para.level or 0)
                
                # Check for bullet points
                if para.level and para.level > 0:
                    parts.append(f"{indent}- {para_text}")
                else:
                    parts.append(f"{para_text}")
        
        return "\n".join(parts)
    
    def _extract_notes(self, slide: Slide) -> Optional[str]:
        """Extract speaker notes from slide.
        
        This is CRITICAL as speaker notes often contain:
        - Key explanations not on slides
        - Presenter's narrative flow
        - Important context and examples
        - Q&A preparation
        
        Args:
            slide: Slide object.
            
        Returns:
            Speaker notes text or None.
        """
        try:
            # Access notes slide
            notes_slide = slide.notes_slide
            if notes_slide is None:
                return None
            
            # Get notes text frame
            if notes_slide.notes_text_frame:
                notes_text = notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    return notes_text
            
            # Fallback: try to get from shapes
            for shape in notes_slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        return text
                        
        except Exception as e:
            logger.debug(f"Could not extract notes: {e}")
        
        return None
    
    def _extract_image(
        self, 
        shape: BaseShape, 
        doc_hash: str, 
        slide_idx: int,
        image_idx: int
    ) -> Optional[Dict[str, Any]]:
        """Extract an image from shape.
        
        Args:
            shape: Shape object (picture).
            doc_hash: Document hash.
            slide_idx: Slide index.
            image_idx: Image index.
            
        Returns:
            Image metadata or None.
        """
        try:
            # Get image data
            image = shape.image
            image_bytes = image.blob
            image_ext = image.ext
            
            # Create storage directory
            image_dir = self.image_storage_dir / doc_hash
            image_dir.mkdir(parents=True, exist_ok=True)
            
            # Generate image ID
            image_id = f"{doc_hash[:8]}_s{slide_idx + 1}_img_{image_idx}"
            image_filename = f"{image_id}.{image_ext}"
            image_path = image_dir / image_filename
            
            # Save image
            with open(image_path, 'wb') as f:
                f.write(image_bytes)
            
            # Get dimensions
            width = shape.width.emu if hasattr(shape, 'width') else 0
            height = shape.height.emu if hasattr(shape, 'height') else 0
            
            return {
                'id': image_id,
                'path': str(image_path),
                'slide': slide_idx + 1,
                'dimensions': {
                    'width': width,
                    'height': height
                }
            }
            
        except Exception as e:
            logger.warning(f"Failed to extract image from slide {slide_idx + 1}: {e}")
            return None
    
    def _table_to_markdown(self, table) -> str:
        """Convert a table to Markdown.
        
        Args:
            table: Table shape object.
            
        Returns:
            Markdown-formatted table.
        """
        markdown_lines = []
        
        for row_idx, row in enumerate(table.rows):
            cells = []
            for cell in row.cells:
                cell_text = cell.text.strip().replace('\n', ' ')
                cells.append(cell_text)
            
            line = "| " + " | ".join(cells) + " |"
            markdown_lines.append(line)
            
            # Add header separator
            if row_idx == 0:
                separator = "| " + " | ".join(["---"] * len(cells)) + " |"
                markdown_lines.append(separator)
        
        return "\n".join(markdown_lines)
    
    def _compute_file_hash(self, file_path: Path) -> str:
        """Compute SHA256 hash of file content."""
        sha256 = hashlib.sha256()
        with open(file_path, 'rb') as f:
            for chunk in iter(lambda: f.read(8192), b''):
                sha256.update(chunk)
        return sha256.hexdigest()
    
    def _extract_title(self, text: str) -> Optional[str]:
        """Extract title from content."""
        lines = text.split('\n')
        
        for line in lines[:20]:
            line = line.strip()
            if line.startswith('# '):
                return line[2:].strip()
        
        for line in lines[:10]:
            line = line.strip()
            if line:
                return line[:100]
        
        return None